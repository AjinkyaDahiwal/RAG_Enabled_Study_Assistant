#extract tables using img2table
from img2table.document import Image as Img2TableImage
import os

from PIL import Image, ImageEnhance, ImageFilter

def preprocess_image_for_table(image_path):
    """
    Preprocess image to enhance contrast and sharpness.
    Returns path to processed image.
    """
    img = Image.open(image_path)
    img = img.convert('L')  # Grayscale
    enhancer = ImageEnhance.Contrast(img)
    img = enhancer.enhance(2)  # Increase contrast
    img = img.filter(ImageFilter.SHARPEN)
    processed_path = image_path.replace('.png', '_cleaned.png').replace('.jpg', '_cleaned.jpg')
    img.save(processed_path)
    return processed_path


def extract_tables_from_image(image_path):
    processed_path = preprocess_image_for_table(image_path)
    """
    Extract tables from an image using img2table library.
    Returns list of table data.
    """
    try:
        doc = Img2TableImage(processed_path)
        tables = doc.extract_tables()
        
        table_data = []
        for table_idx, table in enumerate(tables):
            df = table.df  # Convert table to pandas DataFrame
            table_data.append({
                "table_num": table_idx + 1,
                "dataframe": df,
                "csv": df.to_csv(index=False)  # CSV format for storage
            })
        
        return table_data
    except Exception as e:
        print(f"Error extracting tables: {e}")
        return []

def extract_tables_from_pptx_slides(file_path):
    """
    Extract tables embedded in PowerPoint slides as images.
    """
    from pptx import Presentation
    import tempfile
    
    prs = Presentation(file_path)
    all_tables = []
    
    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Check if shape has an image
            if shape.shape_type == 13:  # Picture shape
                # Save image temporarily
                image = shape.image
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    tmp.write(image.blob)
                    tmp_path = tmp.name
                
                # Extract tables from image
                tables = extract_tables_from_image(tmp_path)
                for table in tables:
                    table["slide_num"] = slide_num + 1
                    all_tables.append(table)
                
                # Clean up
                os.remove(tmp_path)
    
    return all_tables
