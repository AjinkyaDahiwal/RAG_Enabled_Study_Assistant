#extract text from PPTX slides
from pptx import Presentation
from pptx.util import Inches, Pt
import os

def extract_pptx_text(file_path):
    """
    Extract text from all slides in a PowerPoint presentation.
    Returns list of dicts with slide number, title, and text content.
    """
    prs = Presentation(file_path)
    slides_data = []
    
    for slide_num, slide in enumerate(prs.slides):
        slide_info = {
            "slide_num": slide_num + 1,
            "title": "",
            "text": [],
            "shapes": []
        }
        
        # Extract text from all shapes in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # Check if it's the title (usually first shape)
                    if shape == slide.shapes[0] and slide_num == 0:
                        slide_info["title"] = text
                    slide_info["text"].append(text)
        
        # Combine all text from the slide
        slide_info["full_text"] = "\n".join(slide_info["text"])
        slides_data.append(slide_info)
    
    return slides_data


def extract_pptx_metadata(file_path):
    """
    Extract presentation-level metadata.
    """
    prs = Presentation(file_path)
    return {
        "num_slides": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height
    }
