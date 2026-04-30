import os
import sys
import os
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'src')))
from pptx_utils import extract_pptx_text, extract_pptx_metadata
from ocr_utils import extract_text_from_image
from table_utils import extract_tables_from_image,extract_tables_from_pptx_slides,preprocess_image_for_table
from text_preprocessing import normalize_text, clean_text
from pptx_pipeline import process_pptx

def test_extract_pptx_text():
    sample_pptx = "tests/samples/sample2.pptx"
    if os.path.exists(sample_pptx):
        slides = extract_pptx_text(sample_pptx)
        assert len(slides) > 0
        assert "slide_num" in slides[0]
        assert "full_text" in slides[0]

def test_extract_pptx_metadata():
    sample_pptx = "tests/samples/sample2.pptx"
    if os.path.exists(sample_pptx):
        meta = extract_pptx_metadata(sample_pptx)
        assert "num_slides" in meta

def test_normalize_text():
    sample_text = "This   is    a   test    with   extra    spaces"
    normalized = normalize_text(sample_text)
    assert "  " not in normalized

def test_clean_text():
    sample_text = "Page 1\n\nThis is important content.\nPage 2"
    cleaned = clean_text(sample_text)
    assert len(cleaned) > 0

def test_process_pptx_pipeline():
    sample_pptx = "tests/samples/sample2.pptx"
    if os.path.exists(sample_pptx):
        result = process_pptx(sample_pptx)
        assert "slides" in result
        assert "metadata" in result


SAMPLE_PPTX = "tests/samples/sample2.pptx"

def test_extract_tables_from_pptx_slide_images():
    if os.path.exists(SAMPLE_PPTX):
        tables = extract_tables_from_pptx_slides(SAMPLE_PPTX)
        assert isinstance(tables, list)
        # Print first found table as CSV for verification
        if tables:
            print("--- Table from PPTX as CSV ---")
            print(tables[0]["csv"])

def test_pptx_image_ocr_on_slide():
    slides = extract_pptx_text(SAMPLE_PPTX)
    # Find images in the first slide and run OCR
    from pptx import Presentation
    prs = Presentation(SAMPLE_PPTX)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                # Save image temporarily
                import tempfile
                image = shape.image
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    tmp.write(image.blob)
                    tmp_path = tmp.name
                text = extract_text_from_image(tmp_path)
                print(f"--- OCR extracted text from image ---\n{text}\n")
                os.remove(tmp_path)
                # Just test that we can get a non-empty string for image OCR
                assert isinstance(text, str)
                # (Optional) Uncomment to require any text for OCR images:
                # assert len(text) > 0

def test_process_pptx_pipeline_tables_images():
    # End-to-end: does pipeline return tables and process images for OCR?
    if os.path.exists(SAMPLE_PPTX):
        result = process_pptx(SAMPLE_PPTX)
        assert "tables" in result
        assert isinstance(result["tables"], list)
        for table in result["tables"]:
            print("--- Table extracted as CSV ---")
            print(table.get("csv"))
        # Check slides text
        assert "slides" in result
        for slide in result["slides"]:
            print(f"Slide {slide['slide_num']} Title: {slide['title']}")
            print(f"Cleaned Text: ")
            print(slide["cleaned_text"])


def test_ocr_on_sample_image():
    image_path = "tests/samples/sampleimg1.png"
    if os.path.exists(image_path):
        text = extract_text_from_image(image_path)
        print("OCR Text:", text)
        assert isinstance(text, str)
        # (Optional) assert len(text) > 0

def test_table_extraction_on_image():
    image_path = "tests/samples/tableimage2.jpg"
    processed_path = preprocess_image_for_table(image_path)
    if os.path.exists(image_path):
        tables = extract_tables_from_image(processed_path)
        assert isinstance(tables, list)
        if tables:
            for table in tables:
                print("Extracted Table CSV:")
                print(table["csv"])
                print("Extracted Table DataFrame:")
                print(table["dataframe"].to_string())  # Display as raw table, if available
            #print("Extracted Table CSV:", tables[0]["csv"])
         # Add OCR fallback below table extraction
        from ocr_utils import ocr_text_from_image
        text = ocr_text_from_image(processed_path)
        print("OCR Fallback Text:", text)
            

